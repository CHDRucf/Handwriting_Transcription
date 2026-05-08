"""
HTR Senior Project Pitch Deck Generator
Requires: pip install python-pptx
Run:      python3 build_htr_deck.py
Output:   htr_pitch_deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1A, 0x1F, 0x3C)
AMBER       = RGBColor(0xE8, 0xA8, 0x30)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xD4, 0xD8, 0xE8)
CARD_BG     = RGBColor(0x25, 0x2B, 0x52)
MID_GRAY    = RGBColor(0x8A, 0x92, 0xB8)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK_LAYOUT = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────────

def new_slide():
    slide = prs.slides.add_slide(BLANK_LAYOUT)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    return slide


def add_textbox(slide, text, left, top, width, height,
                font_size=24, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_size=20, color=WHITE, line_spacing_pt=6, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(line_spacing_pt)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def heading(slide, text, top=Inches(0.45), font_size=34):
    add_textbox(slide, text, Inches(0.6), top, Inches(12.0), Inches(0.75),
                font_size=font_size, bold=True, color=AMBER)


def amber_bar(slide):
    add_rect(slide, Inches(0.6), Inches(1.22), Inches(12.1), Inches(0.04), AMBER)


def slide_number(slide, n, total=15):
    add_textbox(slide, f"{n} / {total}",
                Inches(11.8), Inches(7.1), Inches(1.4), Inches(0.3),
                font_size=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, 0, 0, W, H, RGBColor(0x0D, 0x11, 0x24))
add_rect(s, 0, 0, W, Inches(0.8), RGBColor(0x12, 0x17, 0x30))

add_textbox(s, "Decoding the Past",
            Inches(0.7), Inches(1.8), Inches(11.9), Inches(1.6),
            font_size=62, bold=True, color=AMBER)
add_textbox(s, "Building an AI-Powered Handwriting Transcription System",
            Inches(0.7), Inches(3.5), Inches(11.0), Inches(1.0),
            font_size=27, color=WHITE)
add_rect(s, Inches(0.7), Inches(4.65), Inches(4.0), Inches(0.05), AMBER)
add_textbox(s, "CS Senior Capstone Project Pitch",
            Inches(0.7), Inches(4.85), Inches(8.0), Inches(0.5),
            font_size=18, color=LIGHT_GRAY)

add_speaker_notes(s,
    "Welcome everyone. I want to show you a project that sits at the intersection of AI, history, "
    "and real-world impact — and pitch it to you as a senior capstone worth your time.")
slide_number(s, 1)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Hook
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "Millions of documents. Completely unreadable by computers.")
amber_bar(s)

add_rect(s, Inches(0.6), Inches(1.35), Inches(5.5), Inches(5.4), CARD_BG)
add_textbox(s, "Handwritten manuscript",
            Inches(0.75), Inches(1.45), Inches(5.2), Inches(0.45),
            font_size=13, color=MID_GRAY)
add_textbox(s, "Flowing cursive, faded ink,\nunsearchable, unindexed.",
            Inches(0.85), Inches(2.1), Inches(5.0), Inches(1.5),
            font_size=21, color=LIGHT_GRAY, italic=True)
add_textbox(s, "[Insert manuscript image here]",
            Inches(0.85), Inches(3.8), Inches(5.0), Inches(0.5),
            font_size=13, color=MID_GRAY, italic=True)

add_textbox(s, "\u2192",
            Inches(6.3), Inches(3.2), Inches(0.9), Inches(0.7),
            font_size=48, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

add_rect(s, Inches(7.3), Inches(1.35), Inches(5.5), Inches(5.4), CARD_BG)
add_textbox(s, "Machine-readable text",
            Inches(7.45), Inches(1.45), Inches(5.2), Inches(0.45),
            font_size=13, color=MID_GRAY)
add_multiline(s,
    ["Dear Sir,", "",
     "I write to you from Philadelphia,",
     "having arrived on the 14th inst.",
     "The weather has been most foul...",
     "", "— searchable, analyzable, usable"],
    Inches(7.5), Inches(2.05), Inches(5.0), Inches(4.2),
    font_size=17, color=WHITE)

add_speaker_notes(s,
    "Here's the problem. The overwhelming majority of historical documents have never been digitized "
    "in any useful way. A computer can't search a handwritten letter. It can't run statistics on a "
    "19th-century ledger. Historians who want to study these materials either read them by hand — "
    "slowly — or just don't. This project is about fixing that.")
slide_number(s, 2)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Why It Matters
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "History from the ground up")
amber_bar(s)

card_data = [
    ("📊", "Distant Reading\nat Scale",
     "Run topic modeling and network\nanalysis across thousands of\ndocuments at once"),
    ("🕵️", "Follow the\nActors",
     "Who traded with whom? Who\nappears in court records,\nand how often?"),
    ("🗣️", "Amplify Unheard\nVoices",
     "Make underrepresented\ncommunities visible in\nthe historical record"),
]
for i, (icon, title, body) in enumerate(card_data):
    x = Inches(0.6 + i * 4.22)
    add_rect(s, x, Inches(1.45), Inches(3.95), Inches(5.3), CARD_BG)
    add_textbox(s, icon, x + Inches(0.2), Inches(1.65),
                Inches(1.0), Inches(0.8), font_size=36)
    add_textbox(s, title, x + Inches(0.2), Inches(2.6),
                Inches(3.5), Inches(0.85), font_size=21, bold=True, color=AMBER)
    add_textbox(s, body, x + Inches(0.2), Inches(3.55),
                Inches(3.5), Inches(2.6), font_size=17, color=LIGHT_GRAY)

add_speaker_notes(s,
    "This isn't just about efficiency. When you make handwritten documents machine-readable, you open "
    "up entire new research questions. Who was trading with whom in 1780? What diseases were spreading "
    "through a village in 1850? Whose names keep appearing in court records? These are questions that "
    "can now be answered computationally — but only if someone builds the tools.")
slide_number(s, 3)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — The Technical Challenge
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "This is harder than it looks")
amber_bar(s)

challenges = [
    ("💧", "Faded or damaged paper",
     "Environmental degradation, water stains, and torn pages confuse OCR systems"),
    ("✂️", "Abbreviations & non-standard spelling",
     "Scribes abbreviated freely; consistent spelling didn't exist in many eras"),
    ("📝", "Marginalia & unusual layouts",
     "Notes in margins, crossed-out words, and non-linear text break standard pipelines"),
    ("✍️", "Wildly varying handwriting styles",
     "Every writer, era, and language has its own style — models can't assume consistency"),
]
for i, (icon, title, body) in enumerate(challenges):
    y = Inches(1.42 + i * 1.38)
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.2), CARD_BG)
    add_textbox(s, icon, Inches(0.8), y + Inches(0.22),
                Inches(0.65), Inches(0.7), font_size=26)
    add_textbox(s, title, Inches(1.65), y + Inches(0.1),
                Inches(4.0), Inches(0.5), font_size=17, bold=True, color=AMBER)
    add_textbox(s, body, Inches(1.65), y + Inches(0.58),
                Inches(10.4), Inches(0.52), font_size=15, color=LIGHT_GRAY)

add_speaker_notes(s,
    "Modern OCR — the technology that reads printed text on your phone — falls apart here. Historical "
    "handwriting is messy, degraded, inconsistent, and often in languages that AI hasn't been trained "
    "on well. The field has been working on this for years, and the state of the art is still far from solved.")
slide_number(s, 4)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Landscape
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "There are tools — but no complete solution")
amber_bar(s)

tools = [
    ("Transkribus",   "Specialized HTR platform",  AMBER),
    ("eScriptorium", "Open-source HTR editor",     AMBER),
    ("Tesseract",    "Open-source OCR engine",     AMBER),
    ("GPT-4o",       "General frontier LLM",       RGBColor(0x6A, 0xC8, 0xA0)),
    ("Claude",       "General frontier LLM",       RGBColor(0x6A, 0xC8, 0xA0)),
    ("Gemini",       "General frontier LLM",       RGBColor(0x6A, 0xC8, 0xA0)),
]
for i, (name, desc, color) in enumerate(tools):
    col, row = i % 3, i // 3
    x = Inches(0.6 + col * 4.22)
    y = Inches(1.48 + row * 2.15)
    add_rect(s, x, y, Inches(3.95), Inches(1.88), CARD_BG)
    add_textbox(s, name, x + Inches(0.2), y + Inches(0.15),
                Inches(3.5), Inches(0.55), font_size=20, bold=True, color=color)
    add_textbox(s, desc, x + Inches(0.2), y + Inches(0.72),
                Inches(3.5), Inches(0.85), font_size=15, color=LIGHT_GRAY)

add_rect(s, Inches(0.6), Inches(6.02), Inches(12.1), Inches(0.85),
         RGBColor(0x35, 0x2A, 0x10))
add_textbox(s,
    "⚠️  No existing system intelligently combines these, evaluates them, and improves over time.",
    Inches(0.85), Inches(6.1), Inches(11.5), Inches(0.65),
    font_size=17, color=AMBER, bold=True)

add_speaker_notes(s,
    "Tools like Transkribus have been around for a while and they work — but they require a lot of "
    "setup and don't leverage the latest AI. Meanwhile, frontier language models like GPT-4 and Claude "
    "are surprisingly good at reading handwriting from images, but they're not designed for this. "
    "Nobody has built the glue layer — a smart pipeline that lets you combine these, evaluate them, "
    "and improve over time. That's what this project is.")
slide_number(s, 5)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Pipeline Vision
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "A full HTR pipeline — built by you")
amber_bar(s)

pipeline = [
    ("📷", "Image\nInput"),
    ("🤖", "Model\nSelection"),
    ("📝", "Transcription\nOutput"),
    ("🔥", "Confidence\nMap"),
    ("✏️", "Human\nEdit"),
    ("🔁", "Feedback\nLoop"),
]
box_w   = Inches(1.72)
box_h   = Inches(2.3)
gap     = Inches(0.2)
start_x = Inches(0.35)
y       = Inches(2.5)

for i, (icon, label) in enumerate(pipeline):
    x = start_x + i * (box_w + gap)
    add_rect(s, x, y, box_w, box_h, CARD_BG)
    add_textbox(s, icon, x, y + Inches(0.22), box_w, Inches(0.75),
                font_size=30, align=PP_ALIGN.CENTER)
    add_textbox(s, label, x, y + Inches(1.1), box_w, Inches(1.0),
                font_size=15, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(pipeline) - 1:
        ax = x + box_w + Inches(0.02)
        add_textbox(s, "\u2192", ax, y + Inches(0.7), gap + Inches(0.1), Inches(0.6),
                    font_size=20, color=AMBER, align=PP_ALIGN.CENTER)

add_textbox(s,
    "Upload a manuscript  \u2192  get a transcription  \u2192  see where it's uncertain  \u2192  fix it  \u2192  make the system smarter.",
    Inches(0.6), Inches(5.1), Inches(12.1), Inches(0.65),
    font_size=17, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_speaker_notes(s,
    "Here's the vision at a high level. You upload a manuscript image. The system picks the best model "
    "for the job — or lets you compare several. It returns a transcription, and critically, it tells "
    "you where it's uncertain. A human editor can correct those spots. And those corrections feed back "
    "in as training examples for the next run. The system gets smarter as it's used.")
slide_number(s, 6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Core Components
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "What you'd actually build")
amber_bar(s)

components = [
    ("Model Router",        "Choose between LLM APIs, local models, or traditional OCR per document"),
    ("Prompt Builder",      "Few-shot examples, anti-hallucination safeguards, temperature control"),
    ("Accuracy Evaluator",  "CER, WER, and BLEU scores calculated automatically per run"),
    ("Confidence Heat Map", "Visualize character-level uncertainty across the transcription"),
    ("Editing Interface",   "Clean UI for human-in-the-loop correction of uncertain passages"),
    ("Feedback System",     "Corrections become few-shot examples that improve future runs"),
]
for i, (comp, desc) in enumerate(components):
    y = Inches(1.42 + i * 0.94)
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.82), CARD_BG)
    add_textbox(s, comp, Inches(0.82), y + Inches(0.1),
                Inches(3.1), Inches(0.55), font_size=16, bold=True, color=AMBER)
    add_textbox(s, desc, Inches(4.1), y + Inches(0.18),
                Inches(8.4), Inches(0.5), font_size=15, color=LIGHT_GRAY)

add_speaker_notes(s,
    "Let me break down the components. This is a system with distinct modules — great for a team "
    "project, because different people can own different pieces. You've got an ML component, a backend, "
    "a frontend, and an evaluation framework all wrapped into one.")
slide_number(s, 7)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Interesting CS Problems
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "The genuinely interesting CS problems")
amber_bar(s)

quads = [
    ("🧠", "Prompt Engineering",
     "Few-shot learning, anti-hallucination\nprompts, temperature tuning —\nmore nuanced than it sounds"),
    ("📊", "Evaluation Frameworks",
     "How do you measure transcription quality?\nCER, WER, BLEU — and when do they\ndisagree with human judgment?"),
    ("🗺️", "Confidence Mapping",
     "Character-level uncertainty visualization\nusing model logits or sampling —\na genuine systems challenge"),
    ("🔁", "Feedback Loops",
     "How do human corrections improve\nfew-shot performance over time?\nDesigning the loop is the hard part"),
]
positions = [
    (Inches(0.6),  Inches(1.45)),
    (Inches(6.87), Inches(1.45)),
    (Inches(0.6),  Inches(4.3)),
    (Inches(6.87), Inches(4.3)),
]
for (icon, title, body), (x, y) in zip(quads, positions):
    add_rect(s, x, y, Inches(6.0), Inches(2.65), CARD_BG)
    add_textbox(s, icon, x + Inches(0.2), y + Inches(0.2),
                Inches(0.8), Inches(0.75), font_size=30)
    add_textbox(s, title, x + Inches(1.1), y + Inches(0.22),
                Inches(4.6), Inches(0.58), font_size=19, bold=True, color=AMBER)
    add_textbox(s, body, x + Inches(0.2), y + Inches(1.05),
                Inches(5.6), Inches(1.45), font_size=15, color=LIGHT_GRAY)

add_speaker_notes(s,
    "I want to flag some of the genuinely interesting CS problems in here, because this isn't just glue "
    "code. Figuring out how to prompt a language model to not hallucinate archaic characters is a research "
    "problem. Building a heat map that shows character-level confidence requires working with model logits "
    "or sampling strategies. Designing a feedback loop that improves few-shot performance over time is a "
    "systems design challenge. There's real depth here.")
slide_number(s, 8)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Research Findings
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "LLMs vs. specialized OCR: it depends")
amber_bar(s)

add_rect(s, Inches(0.6), Inches(1.48), Inches(12.1), Inches(0.55), AMBER)
headers  = ["Criterion", "Specialized Models", "General LLMs"]
col_xs   = [Inches(0.7), Inches(4.0), Inches(8.45)]
col_ws   = [Inches(3.1), Inches(4.2), Inches(4.2)]
for h, x, w in zip(headers, col_xs, col_ws):
    add_textbox(s, h, x, Inches(1.5), w, Inches(0.5),
                font_size=14, bold=True, color=NAVY)

rows_data = [
    ("English accuracy",        "✅ High",           "✅ High"),
    ("Non-English accuracy",    "✅ Better",          "⚠️  Weaker"),
    ("Setup / prep time",       "❌ High",            "✅ Low"),
    ("Cost per page",           "💲 Moderate",        "💲💲 Higher"),
    ("Flexibility / iteration", "⚠️  Limited",        "✅ High"),
    ("Zero-shot capability",    "❌ Needs training",  "✅ Works out of the box"),
]
row_colors = [CARD_BG, RGBColor(0x20, 0x26, 0x48)]
for i, (crit, spec, llm) in enumerate(rows_data):
    y = Inches(2.08 + i * 0.6)
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(0.56), row_colors[i % 2])
    for text, x, w in zip([crit, spec, llm], col_xs, col_ws):
        add_textbox(s, text, x, y + Inches(0.05), w, Inches(0.46), font_size=14, color=WHITE)

add_rect(s, Inches(0.6), Inches(5.72), Inches(12.1), Inches(0.92), RGBColor(0x1A, 0x2E, 0x1A))
add_textbox(s,
    "📌  Two-shot Claude Sonnet 3.5 on whole-page scans outperformed all tested OCR engines "
    "on Belgian historical records  (Kim et al., 2025)",
    Inches(0.85), Inches(5.8), Inches(11.5), Inches(0.76),
    font_size=15, color=RGBColor(0x6A, 0xC8, 0x7A))

add_speaker_notes(s,
    "The research is actually pretty interesting here. Specialized models like Transkribus are great "
    "when you have a lot of data to train on. But recent papers show that a well-prompted LLM — given "
    "just two example transcriptions — can outperform purpose-built OCR systems on handwritten records. "
    "The tradeoff is cost and consistency. Your system would let researchers make that tradeoff "
    "intelligently based on their needs.")
slide_number(s, 9)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Sub-Projects
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "This project has natural sub-teams")
amber_bar(s)

subteams = [
    ("🖥️", "The Pipeline",
     "Model routing · API integration\nBatch processing\nPrompt engineering",
     AMBER),
    ("🎨", "The Interface",
     "Editing UI · Confidence heat map\nUX design\nVisualization",
     RGBColor(0x6A, 0xC8, 0xA0)),
    ("🧪", "The Evaluator",
     "Benchmark framework\nCER / WER / BLEU scoring\nModel comparison reports",
     RGBColor(0x7A, 0xA8, 0xE8)),
]
for i, (icon, title, body, color) in enumerate(subteams):
    x = Inches(0.6 + i * 4.22)
    add_rect(s, x, Inches(1.48), Inches(3.95), Inches(5.25), CARD_BG)
    add_rect(s, x, Inches(1.48), Inches(3.95), Inches(0.18), color)
    add_textbox(s, icon, x + Inches(0.2), Inches(1.82),
                Inches(1.0), Inches(0.78), font_size=34)
    add_textbox(s, title, x + Inches(0.2), Inches(2.72),
                Inches(3.5), Inches(0.6), font_size=21, bold=True, color=color)
    add_textbox(s, body, x + Inches(0.2), Inches(3.45),
                Inches(3.5), Inches(2.9), font_size=17, color=LIGHT_GRAY)

add_speaker_notes(s,
    "One reason this works well as a capstone is that it naturally breaks into parallel workstreams. "
    "You could have a team focused on the ML pipeline and model integration, another building the web "
    "interface and heat map visualization, and another designing the evaluation and benchmarking "
    "framework. These are loosely coupled — you can make real progress independently and then integrate.")
slide_number(s, 10)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Stretch Goals
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "If you want to go further...")
amber_bar(s)

add_textbox(s, "🚀  Stretch Goals",
            Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.65),
            font_size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

stretch = [
    ("🔬", "Fine-tune a local model",
     "Build a growing corpus of corrected transcriptions and use it to fine-tune "
     "an open-source model. Potentially publication-worthy."),
    ("🌍", "Multilingual support",
     "Most tools default to English. Demonstrating improved accuracy on French, Spanish, "
     "or Latin would address a genuine research gap."),
    ("🤝", "Connect to real archives",
     "Integrate with archival databases or digitization pipelines so historians "
     "can use the tool directly on their own collections."),
]
for i, (icon, title, body) in enumerate(stretch):
    y = Inches(2.3 + i * 1.55)
    add_rect(s, Inches(0.6), y, Inches(12.1), Inches(1.35), CARD_BG)
    add_textbox(s, icon, Inches(0.82), y + Inches(0.32),
                Inches(0.7), Inches(0.65), font_size=26)
    add_textbox(s, title, Inches(1.68), y + Inches(0.1),
                Inches(4.0), Inches(0.5), font_size=18, bold=True, color=AMBER)
    add_textbox(s, body, Inches(1.68), y + Inches(0.6),
                Inches(10.2), Inches(0.65), font_size=15, color=LIGHT_GRAY)

add_speaker_notes(s,
    "The baseline project is already meaty. But if your team is ambitious, there's a clear path to "
    "something publication-worthy. Fine-tuning a small open-source model on a domain-specific corpus, "
    "or demonstrating improved accuracy on non-English documents, would be a genuine research contribution.")
slide_number(s, 11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — The Stack
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "The stack")
amber_bar(s)

layers = [
    ("Frontend",   "Web UI for upload, viewing, and transcription editing",
     RGBColor(0x7A, 0xA8, 0xE8)),
    ("Backend",    "Python · FastAPI · job queue for batch processing",
     RGBColor(0x6A, 0xC8, 0xA0)),
    ("Models",     "HuggingFace local models  +  Claude / GPT-4o / Gemini APIs",
     AMBER),
    ("Evaluation", "Custom benchmark harness · CER / WER / BLEU scoring",
     RGBColor(0xD4, 0x7A, 0xC8)),
    ("Data",       "HTR-United open datasets · growing corrected-transcription corpus",
     RGBColor(0xE8, 0x7A, 0x7A)),
]
for i, (label, desc, color) in enumerate(reversed(layers)):
    y = Inches(1.42 + i * 1.02)
    h = Inches(0.9)
    add_rect(s, Inches(0.6),  y, Inches(0.18), h, color)
    add_rect(s, Inches(0.78), y, Inches(11.93), h, CARD_BG)
    add_textbox(s, label, Inches(1.0), y + Inches(0.08),
                Inches(2.2), Inches(0.45), font_size=17, bold=True, color=color)
    add_textbox(s, desc,  Inches(3.4), y + Inches(0.2),
                Inches(9.1), Inches(0.46), font_size=15, color=LIGHT_GRAY)

add_speaker_notes(s,
    "In terms of what you'd actually use: the dataset layer is covered — there are open benchmark "
    "datasets for historical handwriting. The model layer is a mix of API calls to frontier models "
    "and optional local open-source models. The rest is a fairly standard web stack. Python backend, "
    "web frontend. Nothing exotic — but the integration is genuinely interesting.")
slide_number(s, 12)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Why Now
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, 0, 0, W, H, RGBColor(0x0D, 0x11, 0x24))

add_textbox(s, "The tools exist.",
            Inches(0.6), Inches(1.9), Inches(12.1), Inches(0.9),
            font_size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, Inches(3.5), Inches(2.88), Inches(6.3), Inches(0.06), AMBER)

add_textbox(s, "The need is real.",
            Inches(0.6), Inches(3.05), Inches(12.1), Inches(0.9),
            font_size=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, Inches(3.5), Inches(4.03), Inches(6.3), Inches(0.06), AMBER)

add_textbox(s, "Nobody has built the system.",
            Inches(0.6), Inches(4.2), Inches(12.1), Inches(0.9),
            font_size=46, bold=True, color=AMBER, align=PP_ALIGN.CENTER)

add_speaker_notes(s,
    "I want to leave you with this. The individual pieces — the models, the datasets, the evaluation "
    "metrics — they're all out there. The research community has been documenting exactly what works "
    "and what doesn't. What doesn't exist yet is a coherent, usable system that ties them together. "
    "That's the gap this project fills. And it's exactly the kind of thing a motivated senior team can build.")
slide_number(s, 13)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — What You Walk Away With
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
heading(s, "What you'd walk away with")
amber_bar(s)

left_items = [
    ("🧠", "Multimodal LLM integration"),
    ("📐", "Evaluation framework design"),
    ("🔄", "Human-in-the-loop ML systems"),
    ("🌐", "Full-stack web development"),
]
right_items = [
    ("🛠️", "A working tool used by real researchers"),
    ("📄", "A benchmarking study (potentially publishable)"),
    ("💻", "An open-source contribution"),
    ("🤝", "Interdisciplinary collaboration experience"),
]

add_textbox(s, "Technical Skills",
            Inches(0.7), Inches(1.48), Inches(5.5), Inches(0.52),
            font_size=19, bold=True, color=AMBER)
for i, (icon, text) in enumerate(left_items):
    y = Inches(2.05 + i * 1.12)
    add_rect(s, Inches(0.7), y, Inches(5.75), Inches(0.9), CARD_BG)
    add_textbox(s, icon, Inches(0.88), y + Inches(0.14),
                Inches(0.5), Inches(0.6), font_size=22)
    add_textbox(s, text, Inches(1.55), y + Inches(0.2),
                Inches(4.7), Inches(0.5), font_size=17, color=WHITE)

add_textbox(s, "Portfolio Value",
            Inches(7.05), Inches(1.48), Inches(5.7), Inches(0.52),
            font_size=19, bold=True, color=RGBColor(0x6A, 0xC8, 0xA0))
for i, (icon, text) in enumerate(right_items):
    y = Inches(2.05 + i * 1.12)
    add_rect(s, Inches(7.05), y, Inches(5.75), Inches(0.9), CARD_BG)
    add_textbox(s, icon, Inches(7.22), y + Inches(0.14),
                Inches(0.5), Inches(0.6), font_size=22)
    add_textbox(s, text, Inches(7.9), y + Inches(0.2),
                Inches(4.7), Inches(0.5), font_size=17, color=WHITE)

add_speaker_notes(s,
    "This isn't a toy project. Historians actually need this. If you build it well, it will be used. "
    "You'll have hands-on experience with multimodal AI, evaluation methodology, and system design — "
    "and a project you can actually talk about in interviews. The interdisciplinary angle is a plus too.")
slide_number(s, 14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Questions
# ══════════════════════════════════════════════════════════════════════════════
s = new_slide()
add_rect(s, 0, 0, W, H, RGBColor(0x0D, 0x11, 0x24))

add_textbox(s, "What would you build first?",
            Inches(0.7), Inches(2.2), Inches(11.9), Inches(1.45),
            font_size=50, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
add_rect(s, Inches(2.5), Inches(3.78), Inches(8.3), Inches(0.06), AMBER)
add_textbox(s, "Questions & Discussion",
            Inches(0.7), Inches(4.0), Inches(11.9), Inches(0.68),
            font_size=26, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, "[your name  ·  your email  ·  office hours]",
            Inches(0.7), Inches(5.9), Inches(11.9), Inches(0.5),
            font_size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

add_speaker_notes(s,
    "That's the pitch. I'm happy to go deeper on any piece of this — the technical components, "
    "the research background, or how teams might be structured. What questions do you have?")
slide_number(s, 15)

# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
prs.save("htr_pitch_deck.pptx")
print("Saved: htr_pitch_deck.pptx")
print("Notes:")
print("  - Add a real manuscript image to slide 2 where indicated")
print("  - Update contact info on slide 15")
print("  - Emoji rendering varies by OS — check in PowerPoint before presenting")